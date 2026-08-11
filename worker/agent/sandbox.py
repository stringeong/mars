"""폴더 화이트리스트 강제 (최소 권한 원칙, 요구사항 2.4/2.5).

서버가 어떤 경로를 지시하더라도 여기서 검증을 통과하지 못하면 접근을 거부한다.
"""

from pathlib import Path

MAX_FILE_BYTES = 200_000  # 파일당 읽기 상한 (프롬프트 폭주 방지)
MAX_PDF_BYTES = 20_000_000  # PDF 파싱 전 파일 크기 상한
TEXT_EXTENSIONS = {
    ".txt", ".md", ".rst", ".csv", ".tsv", ".json", ".jsonl", ".xml",
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".env", ".log",
    ".py", ".js", ".jsx", ".ts", ".tsx", ".html", ".htm", ".css", ".scss",
    ".sh", ".bash", ".zsh", ".ps1", ".bat", ".cmd", ".sql", ".r", ".java",
    ".c", ".h", ".cc", ".cpp", ".cxx", ".hpp", ".cs", ".go", ".rs", ".rb",
    ".php", ".swift", ".kt", ".kts", ".lua", ".dart", ".vue", ".svelte", ".rtf",
}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx"}
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | DOCUMENT_EXTENSIONS


class SandboxError(PermissionError):
    pass


def _ocr_pdf(path: Path, display_path: str) -> str:
    """Extract text from image-only PDF pages with local Tesseract OCR."""
    try:
        import fitz
        import pytesseract
        from PIL import Image

        parts: list[str] = []
        size = 0
        with fitz.open(path) as document:
            for page in document:
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                image = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
                text = pytesseract.image_to_string(image, lang="kor+eng")
                if text.strip():
                    parts.append(text)
                    size += len(text)
                if size >= MAX_FILE_BYTES:
                    break
        content = "\n".join(parts)[:MAX_FILE_BYTES]
        if not content.strip():
            raise SandboxError(f"PDF OCR에서 텍스트를 찾지 못했습니다: {display_path}")
        return content
    except SandboxError:
        raise
    except Exception as exc:
        raise SandboxError(f"스캔 PDF OCR에 실패했습니다: {display_path} ({exc})") from exc


def _normalize(folders: list[str]) -> list[Path]:
    result = []
    for f in folders:
        try:
            result.append(Path(f).expanduser().resolve())
        except OSError:
            continue
    return result


def is_allowed(path: str | Path, allowed_folders: list[str]) -> bool:
    try:
        target = Path(path).expanduser().resolve()
    except OSError:
        return False
    for base in _normalize(allowed_folders):
        if target == base or base in target.parents:
            return True
    return False


def read_file(path: str, allowed_folders: list[str]) -> str:
    if not is_allowed(path, allowed_folders):
        raise SandboxError(f"허용된 폴더 밖의 경로입니다: {path}")
    p = Path(path).expanduser().resolve()
    if not p.is_file():
        raise SandboxError(f"파일이 아닙니다: {path}")
    suffix = p.suffix.lower()
    if suffix in DOCUMENT_EXTENSIONS and p.stat().st_size > MAX_PDF_BYTES:
        raise SandboxError(f"문서 파일이 너무 큽니다 (최대 {MAX_PDF_BYTES:,} bytes): {path}")
    if suffix == ".pdf":
        if p.stat().st_size > MAX_PDF_BYTES:
            raise SandboxError(f"PDF 파일이 너무 큽니다 (최대 {MAX_PDF_BYTES:,} bytes): {path}")
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(p))
            parts: list[str] = []
            size = 0
            for page in reader.pages:
                text = page.extract_text() or ""
                parts.append(text)
                size += len(text)
                if size >= MAX_FILE_BYTES:
                    break
            content = "\n".join(parts)[:MAX_FILE_BYTES]
            if not content.strip():
                return _ocr_pdf(p, path)
            return content
        except SandboxError:
            raise
        except Exception as exc:
            raise SandboxError(f"PDF 텍스트를 추출할 수 없습니다: {path}") from exc
    if suffix == ".docx":
        try:
            from docx import Document

            document = Document(str(p))
            parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
            for table in document.tables:
                for row in table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
            return _document_content(parts, path)
        except SandboxError:
            raise
        except Exception as exc:
            raise SandboxError(f"Word 문서를 읽을 수 없습니다: {path}") from exc
    if suffix == ".xlsx":
        try:
            from openpyxl import load_workbook

            workbook = load_workbook(str(p), read_only=True, data_only=True)
            parts: list[str] = []
            for worksheet in workbook.worksheets[:10]:
                parts.append(f"[시트: {worksheet.title}]")
                for row in worksheet.iter_rows(values_only=True):
                    values = [str(value) for value in row if value is not None]
                    if values:
                        parts.append(" | ".join(values))
                    if sum(len(part) for part in parts) >= MAX_FILE_BYTES:
                        break
            return _document_content(parts, path)
        except SandboxError:
            raise
        except Exception as exc:
            raise SandboxError(f"Excel 문서를 읽을 수 없습니다: {path}") from exc
    if suffix == ".pptx":
        try:
            from pptx import Presentation

            presentation = Presentation(str(p))
            parts = []
            for index, slide in enumerate(presentation.slides, start=1):
                parts.append(f"[슬라이드 {index}]")
                parts.extend(
                    shape.text for shape in slide.shapes
                    if getattr(shape, "has_text_frame", False) and shape.text
                )
            return _document_content(parts, path)
        except SandboxError:
            raise
        except Exception as exc:
            raise SandboxError(f"PowerPoint 문서를 읽을 수 없습니다: {path}") from exc
    data = p.read_bytes()[:MAX_FILE_BYTES]
    return data.decode("utf-8", errors="replace")


def _document_content(parts: list[str], path: str) -> str:
    content = "\n".join(parts)[:MAX_FILE_BYTES]
    if not content.strip():
        raise SandboxError(f"문서에 추출 가능한 텍스트가 없습니다: {path}")
    return content


# 홈 폴더처럼 큰 경로가 허용되어도 순회 폭주하지 않도록 건너뛰는 디렉터리
SKIP_DIRS = {
    ".git", "node_modules", ".venv", "venv", "__pycache__",
    "Library", "Applications", ".Trash", ".cache", "AppData",
}
MAX_DIRS_VISITED = 2_000  # 순회할 디렉터리 수 상한 (안전장치)


def list_files(allowed_folders: list[str], max_entries: int = 200) -> list[str]:
    """허용 폴더 내 텍스트 파일 목록 (에이전트 컨텍스트 제공용).

    rglob 전체 순회 대신 디렉터리 단위로 돌면서 상한에 도달하면 즉시 멈춘다 —
    사용자가 홈 폴더처럼 큰 경로를 허용해도 워커가 몇 분씩 멈추지 않는다.
    """
    entries: list[str] = []
    visited = 0
    for base in _normalize(allowed_folders):
        if not base.is_dir():
            continue
        stack = [base]
        while stack:
            if len(entries) >= max_entries or visited >= MAX_DIRS_VISITED:
                return entries
            current = stack.pop()
            visited += 1
            try:
                children = sorted(current.iterdir())
            except (PermissionError, OSError):
                continue
            for p in children:
                if p.name.startswith(".") or p.name in SKIP_DIRS:
                    continue
                if p.is_dir():
                    stack.append(p)
                elif p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS:
                    entries.append(str(p))
                    if len(entries) >= max_entries:
                        return entries
    return entries
